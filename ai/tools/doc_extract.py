"""
ai/tools/doc_extract.py
=======================
Deterministic document extraction (no LLM, no model weights). Turns each uploaded
file into a unified set of Chunks tagged with source + location, ready for
session-scoped BM25 retrieval. Tables are kept WHOLE (never split) and rendered
as markdown so numeric ground truth stays intact.

Supported: .pdf (pdfplumber), .docx (python-docx), .pptx (python-pptx),
.xlsx/.xlsm (openpyxl), .csv (stdlib).

This is INFRASTRUCTURE, not a BaseAgent -- extraction needs no reasoning.
"""
from __future__ import annotations

import csv as _csv
from dataclasses import dataclass, field
from pathlib import Path

from ai.src.logger import get_logger

logger = get_logger(__name__)

# RAG-standard chunking: ~500-token windows with ~50-token overlap so a fact split
# across a boundary isn't lost. Token count approximated from words (~1.3 tok/word)
# to avoid a tokenizer dependency. Tables are never windowed (kept whole).
_CHUNK_TOKENS = 500
_OVERLAP_TOKENS = 50
_TOK_PER_WORD = 1.3
_WORDS_PER_CHUNK = int(_CHUNK_TOKENS / _TOK_PER_WORD)      # ~385 words
_OVERLAP_WORDS = int(_OVERLAP_TOKENS / _TOK_PER_WORD)      # ~38 words


def approx_tokens(text: str) -> int:
    return max(1, round(len(text.split()) * _TOK_PER_WORD))


@dataclass
class Chunk:
    text: str
    source: str            # filename
    location: str          # "p.3", "Sheet1", "slide 2", "para"
    is_table: bool = False


@dataclass
class ExtractedDoc:
    source: str
    doc_type: str          # pdf | docx | pptx | xlsx | csv
    chunks: list[Chunk] = field(default_factory=list)
    headings: list[str] = field(default_factory=list)   # for the corpus map


# ── helpers ───────────────────────────────────────────────────────────────────
def _rows_to_markdown(rows: list[list[str]]) -> str:
    rows = [[("" if c is None else str(c)).strip() for c in r] for r in rows if r is not None]
    rows = [r for r in rows if any(cell for cell in r)]
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    out = ["| " + " | ".join(rows[0]) + " |",
           "| " + " | ".join(["---"] * width) + " |"]
    for r in rows[1:]:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


def _group_paragraphs(paras: list[str], source: str, location: str) -> list[Chunk]:
    """Join non-empty paragraphs, then slide a ~500-token window with ~50-token
    overlap across the words. RAG-standard sizing; tables are handled separately."""
    text = "\n".join(p.strip() for p in paras if p and p.strip())
    if not text.strip():
        return []
    words = text.split()
    if len(words) <= _WORDS_PER_CHUNK:
        return [Chunk(text, source, location)]
    chunks, step, i = [], _WORDS_PER_CHUNK - _OVERLAP_WORDS, 0
    while i < len(words):
        piece = words[i:i + _WORDS_PER_CHUNK]
        chunks.append(Chunk(" ".join(piece), source, location))
        if i + _WORDS_PER_CHUNK >= len(words):
            break
        i += step
    return chunks


# ── per-type extractors ───────────────────────────────────────────────────────
def _extract_pdf(path: Path) -> ExtractedDoc:
    import pdfplumber
    doc = ExtractedDoc(path.name, "pdf")
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            loc = f"p.{i}"
            for table in (page.extract_tables() or []):
                md = _rows_to_markdown(table)
                if md:
                    doc.chunks.append(Chunk(md, path.name, loc, is_table=True))
            text = page.extract_text() or ""
            paras = [ln for ln in text.split("\n")]
            doc.chunks.extend(_group_paragraphs(paras, path.name, loc))
    return doc


def _extract_docx(path: Path) -> ExtractedDoc:
    import docx
    d = docx.Document(str(path))
    doc = ExtractedDoc(path.name, "docx")
    paras = []
    for p in d.paragraphs:
        if p.style and p.style.name and p.style.name.lower().startswith("heading") and p.text.strip():
            doc.headings.append(p.text.strip())
        paras.append(p.text)
    doc.chunks.extend(_group_paragraphs(paras, path.name, "body"))
    for ti, table in enumerate(d.tables, start=1):
        rows = [[c.text for c in row.cells] for row in table.rows]
        md = _rows_to_markdown(rows)
        if md:
            doc.chunks.append(Chunk(md, path.name, f"table {ti}", is_table=True))
    return doc


def _extract_pptx(path: Path) -> ExtractedDoc:
    from pptx import Presentation
    prs = Presentation(str(path))
    doc = ExtractedDoc(path.name, "pptx")
    for si, slide in enumerate(prs.slides, start=1):
        loc = f"slide {si}"
        texts = []
        for shape in slide.shapes:
            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                md = _rows_to_markdown(rows)
                if md:
                    doc.chunks.append(Chunk(md, path.name, loc, is_table=True))
            elif shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())
        if texts:
            if texts[0]:
                doc.headings.append(texts[0])
            doc.chunks.extend(_group_paragraphs(texts, path.name, loc))
    return doc


def _extract_xlsx(path: Path) -> ExtractedDoc:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    doc = ExtractedDoc(path.name, "xlsx")
    for ws in wb.worksheets:
        rows = [[c for c in row] for row in ws.iter_rows(values_only=True)]
        md = _rows_to_markdown(rows)
        if md:
            doc.headings.append(f"Sheet: {ws.title}")
            doc.chunks.append(Chunk(md, path.name, f"sheet {ws.title}", is_table=True))
    wb.close()
    return doc


def _extract_csv(path: Path) -> ExtractedDoc:
    doc = ExtractedDoc(path.name, "csv")
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        rows = list(_csv.reader(f))
    md = _rows_to_markdown(rows)
    if md:
        doc.chunks.append(Chunk(md, path.name, "csv", is_table=True))
    return doc


_DISPATCH = {
    ".pdf": _extract_pdf, ".docx": _extract_docx, ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx, ".xlsm": _extract_xlsx, ".csv": _extract_csv,
}


def extract(path: str | Path) -> ExtractedDoc:
    """Extract a single file into an ExtractedDoc. Unsupported/erroring files
    return an empty doc (logged) rather than raising, so one bad upload can't
    sink the whole run."""
    path = Path(path)
    fn = _DISPATCH.get(path.suffix.lower())
    if fn is None:
        logger.warning("doc_extract: unsupported file type %s", path.suffix)
        return ExtractedDoc(path.name, path.suffix.lstrip(".") or "unknown")
    try:
        doc = fn(path)
        logger.info("doc_extract: %s -> %d chunks (%d tables)", path.name,
                    len(doc.chunks), sum(1 for c in doc.chunks if c.is_table))
        return doc
    except Exception as e:
        logger.warning("doc_extract: failed on %s (%s); skipping", path.name, e)
        return ExtractedDoc(path.name, path.suffix.lstrip("."))


def extract_all(paths: list[str | Path]) -> list[ExtractedDoc]:
    return [extract(p) for p in paths]